"""
Gera apresentacoes PPTX para cada conector do framework ELT.
Uma apresentacao por conector, com slides de titulo, configuracao,
exemplo e tratamento de erros.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DOCS_DIR = os.path.join(os.path.dirname(__file__), "..", "docs")
os.makedirs(DOCS_DIR, exist_ok=True)

BLUE = RGBColor(0x1F, 0x49, 0x7D)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF2, 0xF2, 0xF2)


def _add_title_box(slide, text, left=Inches(0.8), top=Inches(0.4),
                   width=Inches(11), height=Inches(0.8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BLUE
    return txBox


def _add_subtitle_box(slide, text, left=Inches(0.8), top=Inches(1.2),
                      width=Inches(11), height=Inches(0.5)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    return txBox


def _add_body_box(slide, lines, left=Inches(0.8), top=Inches(1.9),
                  width=Inches(11), height=Inches(4.5), font_size=Pt(12)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.color.rgb = DARK
        p.space_after = Pt(4)
        if line.startswith("##"):
            p.text = line.replace("## ", "")
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = BLUE
            p.space_before = Pt(12)
        elif line.startswith("#"):
            p.text = line.replace("# ", "")
            p.font.size = Pt(14)
            p.font.bold = True
            p.space_before = Pt(8)
    return txBox


def _add_code_box(slide, code_text, left=Inches(0.8), top=Inches(1.9),
                  width=Inches(11), height=Inches(4.5)):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x2D, 0x2D, 0x2D)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    for i, line in enumerate(code_text.strip().split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD4)
        p.font.name = "Consolas"
        p.space_after = Pt(1)
    return shape


def _add_footer(slide, text="Framework ELT - Generic ELT Pipeline"):
    txBox = slide.shapes.add_textbox(
        Inches(0.8), Inches(6.8), Inches(11), Inches(0.4)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER


def _build_prs():
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    return prs


def _add_slide(prs):
    layout = prs.slide_layouts[6]  # Blank
    return prs.slides.add_slide(layout)


# ============================================================
# 1. ORACLE
# ============================================================
def build_oracle():
    prs = _build_prs()
    # Slide 1 - Titulo
    s = _add_slide(prs)
    _add_title_box(s, "Conector Oracle", top=Inches(2.5))
    _add_subtitle_box(s, "Extracao de dados de bancos Oracle para a camada Bronze",
                      top=Inches(3.3))
    _add_footer(s)

    # Slide 2 - Como funciona
    s = _add_slide(prs)
    _add_title_box(s, "Como Funciona")
    _add_body_box(s, [
        "## Fluxo",
        "Airflow Connection (conn_id) -> OracleConnector -> OracleExtractor -> bronze",
        "",
        "## Arquivos",
        "src/connectors/oracle_connector.py  -  Conector SQLAlchemy + oracledb",
        "src/extractors/oracle.py  -  Extrator BaseExtractor",
        "",
        "## Driver",
        "Requer oracledb (ou cx_Oracle) instalado no Airflow: pip install oracledb",
    ])
    _add_footer(s)

    # Slide 3 - Configuracao
    s = _add_slide(prs)
    _add_title_box(s, "Configuracao")
    _add_body_box(s, [
        "## Coluna config (JSONB) - apenas o connection_airflow",
        '{"connection_airflow": "sisus_oracle"}',
        "",
        "## Colunas da schedule (campos diretos)",
        "schema_source  -  Schema (owner) da tabela no Oracle",
        "table_source  -  Nome da tabela no Oracle",
        "query_source  -  Query SQL customizada (Oracle syntax)",
        "",
        "## Connection no Airflow",
        "Conn Type: oracle | Schema: service_name | Port: 1521",
    ])
    _add_footer(s)

    # Slide 4 - Exemplo SQL
    s = _add_slide(prs)
    _add_title_box(s, "Exemplo: INSERT na Schedule")
    _add_code_box(s, """INSERT INTO global.schedule (
    type_source, layer, projeto, ordem, ativo,
    schema_source, table_source,
    schema_destiny, table_destiny,
    strategy_destiny, schedule_cron, config
) VALUES (
    'ORACLE', 'bronze', 'sisus', 1, true,
    'SISUS', 'TB_PACIENTE',
    'sisus', 'tb_paciente',
    'truncate', '0 3 * * *',
    '{"connection_airflow": "sisus_oracle"}'::jsonb
);""")
    _add_footer(s)

    # Slide 5 - Erros
    s = _add_slide(prs)
    _add_title_box(s, "Tratamento de Erros")
    _add_body_box(s, [
        "## ORA-01843 / ORA-01861 (formato de data)",
        "Extrator refaz a query com casts automaticos para text",
        "",
        "## connection_airflow ausente",
        "Erro antes da extracao, registro marcado como erro na controle_execucao",
        "",
        "## Tabela inexistente",
        "Erro propagado do Oracle, registrado na controle_execucao",
    ])
    _add_footer(s)

    path = os.path.join(DOCS_DIR, "Conector_Oracle.pptx")
    prs.save(path)
    print(f"  -> {path}")


# ============================================================
# 2. POSTGRES
# ============================================================
def build_postgres():
    prs = _build_prs()
    s = _add_slide(prs)
    _add_title_box(s, "Conector PostgreSQL", top=Inches(2.5))
    _add_subtitle_box(s, "Extracao de dados de bancos PostgreSQL para a camada Bronze",
                      top=Inches(3.3))
    _add_footer(s)

    s = _add_slide(prs)
    _add_title_box(s, "Como Funciona")
    _add_body_box(s, [
        "## Fluxo",
        "Airflow Connection (conn_id) -> PostgresConnector -> PostgresExtractor -> bronze",
        "",
        "## Arquivos",
        "src/connectors/postgres_connector.py  -  Conector SQLAlchemy",
        "src/extractors/postgres.py  -  Extrator BaseExtractor",
        "",
        "## Fallback automatico",
        "Se a query falhar com ano fora do range, refaz com todos os cast ::text",
    ])
    _add_footer(s)

    s = _add_slide(prs)
    _add_title_box(s, "Configuracao")
    _add_body_box(s, [
        "## Coluna config (JSONB)",
        '{"connection_airflow": "caju"}',
        "",
        "## Colunas da schedule (campos diretos)",
        "schema_source  -  Schema da tabela (default: public)",
        "table_source  -  Nome da tabela",
        "query_source  -  Query SQL customizada (PostgreSQL syntax)",
        "",
        "## Connection no Airflow",
        "Conn Type: postgres | Schema: database | Port: 5432",
    ])
    _add_footer(s)

    s = _add_slide(prs)
    _add_title_box(s, "Exemplo: INSERT na Schedule")
    _add_code_box(s, """INSERT INTO global.schedule (
    type_source, layer, projeto, ordem, ativo,
    schema_source, table_source,
    schema_destiny, table_destiny,
    strategy_destiny, schedule_cron, config
) VALUES (
    'POSTGRE', 'bronze', 'cnes', 1, true,
    'public', 'estabelecimento',
    'cnes', 'estabelecimento',
    'truncate', '0 4 * * *',
    '{"connection_airflow": "caju"}'::jsonb
);""")
    _add_footer(s)

    path = os.path.join(DOCS_DIR, "Conector_Postgres.pptx")
    prs.save(path)
    print(f"  -> {path}")


# ============================================================
# 3. GOOGLE SHEETS
# ============================================================
def build_google_sheets():
    prs = _build_prs()
    s = _add_slide(prs)
    _add_title_box(s, "Conector Google Sheets", top=Inches(2.5))
    _add_subtitle_box(s, "Extracao de planilhas Google Sheets para a camada Bronze",
                      top=Inches(3.3))
    _add_footer(s)

    s = _add_slide(prs)
    _add_title_box(s, "Como Funciona")
    _add_body_box(s, [
        "## Fluxo",
        "URL da planilha -> gspread API -> GoogleSheetsExtractor -> bronze",
        "",
        "## Arquivos",
        "src/extractors/google_sheets.py  -  Extrator BaseExtractor",
        "",
        "## Autenticacao",
        "Via Airflow Connection (conn_id) ou arquivo credentials.json",
        "Suporta planilhas nativas e arquivos .xlsx no Google Drive",
    ])
    _add_footer(s)

    s = _add_slide(prs)
    _add_title_box(s, "Configuracao")
    _add_body_box(s, [
        "## Coluna config (JSONB) - opcional",
        '{"connection_airflow": "google_sheets"}',
        "Se vazio, usa o arquivo credentials.json local",
        "",
        "## Colunas da schedule (campos diretos)",
        "url  -  URL da planilha Google Sheets",
        "table_source  -  Nome da aba (worksheet)",
        "header_row_source  -  Numero da linha do cabecalho",
        "columns_source  -  Colunas esperadas (separadas por virgula)",
    ])
    _add_footer(s)

    s = _add_slide(prs)
    _add_title_box(s, "Exemplo: INSERT na Schedule")
    _add_code_box(s, """INSERT INTO global.schedule (
    type_source, layer, projeto, ordem, ativo,
    table_source, header_row_source,
    schema_destiny, table_destiny,
    url, strategy_destiny, schedule_cron, config
) VALUES (
    'GOOGLE_SHEETS', 'bronze', 'mais_medicos', 1, true,
    'PERCURSO_PROFISSIONAL', 2,
    'mais_medicos', 'percurso_profissional',
    'https://docs.google.com/spreadsheets/d/XXX/edit',
    'truncate', '0 4 * * *',
    '{"connection_airflow": "google_sheets"}'::jsonb
);""")
    _add_footer(s)

    path = os.path.join(DOCS_DIR, "Conector_Google_Sheets.pptx")
    prs.save(path)
    print(f"  -> {path}")


# ============================================================
# 4. API
# ============================================================
def build_api():
    prs = _build_prs()
    s = _add_slide(prs)
    _add_title_box(s, "Conector API", top=Inches(2.5))
    _add_subtitle_box(s, "Extracao de dados via API REST para a camada Bronze",
                      top=Inches(3.3))
    _add_footer(s)

    s = _add_slide(prs)
    _add_title_box(s, "Como Funciona")
    _add_body_box(s, [
        "## Fluxo",
        "Airflow Connection (conn_id) -> requests -> ApiExtractor -> bronze",
        "",
        "## Arquivos",
        "src/extractors/api.py  -  Extrator BaseExtractor",
        "",
        "## Autenticacao",
        "A Connection do Airflow fornece login/password para Basic Auth",
        "A URL base e montada a partir de base_url + endpoint no config",
    ])
    _add_footer(s)

    s = _add_slide(prs)
    _add_title_box(s, "Configuracao (coluna config)")
    _add_body_box(s, [
        "## Obrigatorio",
        "connection_airflow  -  ID da Connection no Airflow (contem login/senha)",
        "",
        "## Opcional",
        "base_url  -  URL base da API (ex.: https://api.example.com/v1)",
        "endpoint  -  Endpoint (ex.: /pacientes)",
        "method  -  GET, POST, PUT, DELETE (default: GET)",
        "headers  -  Headers adicionais (JSON)",
        "params  -  Query parameters (JSON)",
        "body  -  Corpo da requisicao (JSON, para POST/PUT)",
    ])
    _add_footer(s)

    s = _add_slide(prs)
    _add_title_box(s, "Exemplo: INSERT na Schedule")
    _add_code_box(s, """INSERT INTO global.schedule (
    type_source, layer, projeto, ordem, ativo,
    schema_destiny, table_destiny,
    strategy_destiny, schedule_cron, config
) VALUES (
    'API', 'bronze', 'meu_projeto', 1, true,
    'meu_projeto', 'dados_api',
    'truncate', '0 6 * * *',
    '{
        "connection_airflow": "minha_api",
        "base_url": "https://api.example.com/v1",
        "endpoint": "/pacientes",
        "method": "GET",
        "params": {"uf": "PE", "ano": "2025"}
    }'::jsonb
);""")
    _add_footer(s)

    s = _add_slide(prs)
    _add_title_box(s, "Tratamento de Erros HTTP")
    _add_body_box(s, [
        "## 401 - Autenticacao falhou",
        "Verifique login/senha na Connection do Airflow",
        "",
        "## 403 - Acesso negado",
        "Verifique permissoes do usuario na API",
        "",
        "## 404 - Recurso nao encontrado",
        "Verifique base_url e endpoint no config",
        "",
        "## 5xx - Erro no servidor remoto",
        "Erro propagado, registrado na controle_execucao",
        "",
        "## Respostas suportadas",
        "JSON (list/dict), CSV, arquivo (attachment), conteudo binario",
    ])
    _add_footer(s)

    path = os.path.join(DOCS_DIR, "Conector_API.pptx")
    prs.save(path)
    print(f"  -> {path}")


# ============================================================
# 5. FTP / DATASUS
# ============================================================
def build_ftp():
    prs = _build_prs()
    s = _add_slide(prs)
    _add_title_box(s, "Conector FTP / DATASUS", top=Inches(2.5))
    _add_subtitle_box(s, "Extracao de arquivos .dbc do FTP do DATASUS",
                      top=Inches(3.3))
    _add_footer(s)

    s = _add_slide(prs)
    _add_title_box(s, "Como Funciona")
    _add_body_box(s, [
        "## Fluxo",
        "FTP datasus.gov.br -> download .dbc -> converte .dbf -> DataFrame -> bronze",
        "",
        "## Arquivos",
        "src/extractors/ftp.py  -  Extrator FTPDatasusExtractor",
        "",
        "## Paralelismo",
        "Download paralelo por UF com ThreadPoolExecutor",
        "Cada UF gera um DataFrame que e concatenado no final",
    ])
    _add_footer(s)

    s = _add_slide(prs)
    _add_title_box(s, "Configuracao (coluna config)")
    _add_body_box(s, [
        "## Opcional (tem defaults)",
        "ftp_host  -  Host do FTP (default: ftp.datasus.gov.br)",
        "ftp_base  -  Caminho base (default: dissemin/publicos/CNES/200508_/Dados)",
        "theme  -  Tema/SIGLA (ex.: ST, CNES, SIH)",
        "ano_mes  -  Periodo (ex.: 2604 = abril/2026)",
        "ufs  -  Lista de UFs (default: todas)",
        "max_workers  -  Threads paralelas (default: 5)",
        "encoding  -  Encoding dos arquivos (default: latin-1)",
    ])
    _add_footer(s)

    s = _add_slide(prs)
    _add_title_box(s, "Exemplo: INSERT na Schedule")
    _add_code_box(s, """INSERT INTO global.schedule (
    type_source, layer, projeto, ordem, ativo,
    schema_destiny, table_destiny,
    strategy_destiny, schedule_cron, config
) VALUES (
    'FTP', 'bronze', 'cnes', 1, true,
    'cnes', 'estabelecimento',
    'truncate', '0 2 * * *',
    '{
        "theme": "ST",
        "ano_mes": "2604",
        "ufs": ["PE", "BA", "CE"],
        "max_workers": 5
    }'::jsonb
);""")
    _add_footer(s)

    path = os.path.join(DOCS_DIR, "Conector_FTP_DATASUS.pptx")
    prs.save(path)
    print(f"  -> {path}")


# ============================================================
# 6. S3 / CKAN / CSV_URL
# ============================================================
def build_s3():
    prs = _build_prs()
    s = _add_slide(prs)
    _add_title_box(s, "Conector S3 / CKAN / CSV_URL", top=Inches(2.5))
    _add_subtitle_box(s, "Extracao de arquivos CSV/Parquet de URLs publicas ou S3",
                      top=Inches(3.3))
    _add_footer(s)

    s = _add_slide(prs)
    _add_title_box(s, "Como Funciona")
    _add_body_box(s, [
        "## Fluxo",
        "URL (HTTP/S3) -> download -> parse CSV/Parquet -> DataFrame -> bronze",
        "",
        "## Arquivos",
        "src/extractors/s3.py  -  Extrator S3Extractor",
        "",
        "## Formatos suportados",
        "CSV (com ou sem zip), Parquet, CSV delimitado",
        "Um unico extrator serve para S3, CKAN e CSV_URL",
    ])
    _add_footer(s)

    s = _add_slide(prs)
    _add_title_box(s, "Configuracao (coluna config)")
    _add_body_box(s, [
        "## Opcional (tem defaults)",
        "delimiter  -  Delimitador CSV (default: ;)",
        "encoding  -  Encoding do arquivo (default: latin-1)",
        "compression  -  zip ou null",
        "header_row  -  Numero da linha do cabecalho (default: 0)",
        "",
        "## Coluna url (obrigatoria)",
        "url  -  URL do arquivo para download",
    ])
    _add_footer(s)

    s = _add_slide(prs)
    _add_title_box(s, "Exemplo: INSERT na Schedule")
    _add_code_box(s, """INSERT INTO global.schedule (
    type_source, layer, projeto, ordem, ativo,
    schema_destiny, table_destiny,
    url, strategy_destiny, schedule_cron, config
) VALUES (
    'CSV_URL', 'bronze', 'dados_abertos', 1, true,
    'dados_abertos', 'estabelecimentos',
    'https://dados.gov.br/dados/conjuntos-dados/XXX.csv',
    'truncate', '0 4 * * 1',
    '{"delimiter": ";", "encoding": "latin-1"}'::jsonb
);""")
    _add_footer(s)

    path = os.path.join(DOCS_DIR, "Conector_S3_CKAN_CSV_URL.pptx")
    prs.save(path)
    print(f"  -> {path}")


# ============================================================
# 7. MINIO
# ============================================================
def build_minio():
    prs = _build_prs()
    s = _add_slide(prs)
    _add_title_box(s, "Conector MinIO", top=Inches(2.5))
    _add_subtitle_box(s, "Extracao de arquivos .parquet do MinIO",
                      top=Inches(3.3))
    _add_footer(s)

    s = _add_slide(prs)
    _add_title_box(s, "Como Funciona")
    _add_body_box(s, [
        "## Fluxo",
        "MinIO (S3 compativel) -> download .parquet -> DataFrame -> bronze",
        "",
        "## Arquivos",
        "src/extractors/minio.py  -  Extrator MinioExtractor",
        "src/connectors/minio_connector.py  -  Conector MinIO (upload/download)",
    ])
    _add_footer(s)

    s = _add_slide(prs)
    _add_title_box(s, "Configuracao (coluna config)")
    _add_body_box(s, [
        "## Opcional (tem defaults via env vars)",
        "endpoint  -  Host:porta do MinIO",
        "access_key  -  Chave de acesso",
        "secret_key  -  Chave secreta",
        "secure  -  true/false (default: false)",
        "bucket  -  Nome do bucket",
        "object_name  -  Caminho do objeto no bucket",
        "",
        "## Alternativa: coluna url",
        "url  -  Caminho do objeto (se object_name nao informado)",
    ])
    _add_footer(s)

    s = _add_slide(prs)
    _add_title_box(s, "Exemplo: INSERT na Schedule")
    _add_code_box(s, """INSERT INTO global.schedule (
    type_source, layer, projeto, ordem, ativo,
    schema_destiny, table_destiny,
    strategy_destiny, schedule_cron, config
) VALUES (
    'MINIO', 'bronze', 'cnes', 1, true,
    'cnes', 'estabelecimento',
    'truncate', '0 5 * * *',
    '{
        "endpoint": "minio.example.com:9000",
        "bucket": "cnes-bronze",
        "object_name": "st/2604/data.parquet",
        "secure": false
    }'::jsonb
);""")
    _add_footer(s)

    path = os.path.join(DOCS_DIR, "Conector_MinIO.pptx")
    prs.save(path)
    print(f"  -> {path}")


# ============================================================
# 8. XLSX
# ============================================================
def build_xlsx():
    prs = _build_prs()
    s = _add_slide(prs)
    _add_title_box(s, "Conector XLSX", top=Inches(2.5))
    _add_subtitle_box(s, "Extracao de planilhas Excel (.xlsx) locais",
                      top=Inches(3.3))
    _add_footer(s)

    s = _add_slide(prs)
    _add_title_box(s, "Como Funciona")
    _add_body_box(s, [
        "## Fluxo",
        "Arquivo local .xlsx -> openpyxl -> DataFrame -> bronze",
        "",
        "## Arquivos",
        "src/extractors/xlsx.py  -  Extrator XlsxExtractor",
        "",
        "## Observacao",
        "O arquivo deve existir no filesystem do worker Airflow",
    ])
    _add_footer(s)

    s = _add_slide(prs)
    _add_title_box(s, "Configuracao")
    _add_body_box(s, [
        "## Coluna config (JSONB)",
        '{"file_path": "/caminho/para/arquivo.xlsx"}',
        "",
        "## Colunas da schedule (alternativa)",
        "url  -  Caminho do arquivo (se file_path nao informado no config)",
        "table_source  -  Nome da aba (worksheet)",
        "header_row_source  -  Numero da linha do cabecalho",
    ])
    _add_footer(s)

    s = _add_slide(prs)
    _add_title_box(s, "Exemplo: INSERT na Schedule")
    _add_code_box(s, """INSERT INTO global.schedule (
    type_source, layer, projeto, ordem, ativo,
    table_source, header_row_source,
    schema_destiny, table_destiny,
    url, strategy_destiny, schedule_cron, config
) VALUES (
    'XLSX', 'bronze', 'meu_projeto', 1, true,
    'Planilha1', 1,
    'meu_projeto', 'dados_excel',
    '/data/entrada/relatorio.xlsx',
    'truncate', '0 6 * * *',
    '{"file_path": "/data/entrada/relatorio.xlsx"}'::jsonb
);""")
    _add_footer(s)

    path = os.path.join(DOCS_DIR, "Conector_XLSX.pptx")
    prs.save(path)
    print(f"  -> {path}")


# ============================================================
# MAIN
# ============================================================
if __name__ == "__main__":
    print("Gerando apresentacoes dos conectores...")
    build_oracle()
    build_postgres()
    build_google_sheets()
    build_api()
    build_ftp()
    build_s3()
    build_minio()
    build_xlsx()
    print("Concluido!")
